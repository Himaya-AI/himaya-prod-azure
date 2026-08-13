"""
File → text extraction adapter.

A single, shared adapter that turns any supported file (from SharePoint /
OneDrive / Teams / email attachments) into plain text so it can be sent to the
dlp-classifier service. Previously this logic lived only inside the email-DLP
module (`dlp_inline._extract_text_from_attachment`), so the DSPM scanners
(SharePoint / Teams) only handled PDF + naive utf-8 decode — DOCX/XLSX/PPTX
files were fed to the classifier as raw zip bytes (garbage).

Supported: PDF, plain text (txt/csv/log/json/xml/md/html), DOCX, XLSX, PPTX,
images (OCR), RTF, EML. Every backend is import-guarded and best-effort: an
unsupported or unreadable file returns "" rather than raising.

Use `extract_text_async()` from async code — heavy parsing (PDF/OCR) runs in a
thread executor with a timeout so it never blocks the event loop.
"""
from __future__ import annotations

import asyncio
import io
import logging
import re

logger = logging.getLogger(__name__)

_DEFAULT_MAX_CHARS = 50_000
_DEFAULT_TIMEOUT_SEC = 20.0

_TEXT_EXTS = (".txt", ".csv", ".log", ".json", ".xml", ".md", ".html", ".htm", ".tsv", ".yaml", ".yml", ".ini", ".conf")
_IMAGE_EXTS = (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp")


def _extract_pdf(content: bytes) -> str:
    from pdfminer.high_level import extract_text as pdf_extract_text
    from pdfminer.pdfparser import PDFSyntaxError
    try:
        return pdf_extract_text(io.BytesIO(content), maxpages=20) or ""
    except PDFSyntaxError as e:
        logger.debug(f"file_extractor: PDF syntax error: {e}")
        return ""


def _extract_docx(content: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(content))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n".join(parts)


def _extract_xlsx(content: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    parts = []
    try:
        for sheet in wb.worksheets[:5]:
            for row in sheet.iter_rows(max_row=500, values_only=True):
                row_text = " ".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    parts.append(row_text)
    finally:
        wb.close()
    return "\n".join(parts)


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    parts = []
    for slide in prs.slides[:30]:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
    return "\n".join(parts)


def _extract_image_ocr(content: bytes) -> str:
    from PIL import Image
    import pytesseract
    img = Image.open(io.BytesIO(content))
    max_dim = 2000
    if img.width > max_dim or img.height > max_dim:
        ratio = min(max_dim / img.width, max_dim / img.height)
        img = img.resize((int(img.width * ratio), int(img.height * ratio)))
    return pytesseract.image_to_string(img, timeout=30) or ""


def _extract_rtf(content: bytes) -> str:
    from striprtf.striprtf import rtf_to_text
    return rtf_to_text(content.decode("utf-8", errors="ignore")) or ""


def _extract_eml(content: bytes) -> str:
    import email
    from email import policy
    msg = email.message_from_bytes(content, policy=policy.default)
    parts = [
        f"Subject: {msg.get('subject', '')}",
        f"From: {msg.get('from', '')}",
        f"To: {msg.get('to', '')}",
    ]
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                payload = part.get_payload(decode=True)
                if payload:
                    parts.append(payload.decode("utf-8", errors="ignore"))
    else:
        payload = msg.get_payload(decode=True)
        if payload:
            parts.append(payload.decode("utf-8", errors="ignore"))
    return "\n".join(parts)


def _extract_plaintext(content: bytes, name: str, content_type: str) -> str:
    text = content.decode("utf-8", errors="ignore")
    if name.endswith((".html", ".htm")) or "html" in content_type:
        text = re.sub(r"<[^>]+>", " ", text)
        text = re.sub(r"\s+", " ", text)
    return text


def extract_text(
    content: bytes,
    filename: str = "",
    content_type: str = "",
    *,
    max_chars: int = _DEFAULT_MAX_CHARS,
) -> str:
    """Synchronously extract plain text from file bytes.

    Returns "" (never raises) if the type is unsupported or extraction fails.
    Prefer `extract_text_async` from async code.
    """
    if not content:
        return ""
    name = (filename or "").lower()
    ctype = (content_type or "").lower()

    def _pick():
        if ctype == "application/pdf" or name.endswith(".pdf"):
            return _extract_pdf
        if ctype.startswith("text/") or name.endswith(_TEXT_EXTS):
            return lambda c: _extract_plaintext(c, name, ctype)
        if name.endswith(".docx") or "wordprocessingml.document" in ctype:
            return _extract_docx
        if name.endswith(".xlsx") or "spreadsheetml.sheet" in ctype:
            return _extract_xlsx
        if name.endswith(".pptx") or "presentationml.presentation" in ctype:
            return _extract_pptx
        if ctype.startswith("image/") or name.endswith(_IMAGE_EXTS):
            return _extract_image_ocr
        if ctype == "application/rtf" or name.endswith(".rtf"):
            return _extract_rtf
        if ctype == "message/rfc822" or name.endswith(".eml"):
            return _extract_eml
        return None

    handler = _pick()
    if handler is None:
        # Last resort: if it decodes to mostly-printable text, use it; else skip.
        try:
            text = content.decode("utf-8")
        except Exception:
            return ""
        printable = sum(c.isprintable() or c.isspace() for c in text[:2000])
        if text and printable / max(len(text[:2000]), 1) > 0.85:
            return text[:max_chars]
        return ""

    try:
        text = handler(content)
    except ImportError as e:
        logger.debug(f"file_extractor: parser unavailable for {name or ctype}: {e}")
        return ""
    except Exception as e:
        logger.warning(f"file_extractor: extraction failed for {name or ctype}: {e}")
        return ""
    return (text or "")[:max_chars]


async def extract_text_async(
    content: bytes,
    filename: str = "",
    content_type: str = "",
    *,
    max_chars: int = _DEFAULT_MAX_CHARS,
    timeout: float = _DEFAULT_TIMEOUT_SEC,
) -> str:
    """Async wrapper — runs `extract_text` in a thread executor with a timeout."""
    loop = asyncio.get_event_loop()
    try:
        return await asyncio.wait_for(
            loop.run_in_executor(
                None, lambda: extract_text(content, filename, content_type, max_chars=max_chars)
            ),
            timeout=timeout,
        )
    except asyncio.TimeoutError:
        logger.warning(f"file_extractor: extraction timed out for {filename or content_type}")
        return ""
    except Exception as e:
        logger.warning(f"file_extractor: async extraction error for {filename or content_type}: {e}")
        return ""
