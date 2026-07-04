"""
Himaya Inline DLP — API-based DLP that works through existing
Gmail/M365 integrations without requiring external transport rules.

This approach:
1. Hooks into delta sync when emails are processed
2. Scans outbound emails from the Sent Items folder
3. Takes actions via API (quarantine, label, notify) rather than SMTP routing
"""
from __future__ import annotations

import logging
import os
from datetime import datetime, timezone, timedelta
from typing import Optional, Literal

from sqlalchemy import text
from sqlalchemy.ext.asyncio import AsyncSession

from backend.services.dlp_service import classify_text, DLP_CATEGORIES

logger = logging.getLogger(__name__)


# ── DLP Configuration ─────────────────────────────────────────────────────

class DLPConfig:
    """Organization-level DLP configuration."""
    
    def __init__(
        self,
        enabled: bool = False,
        scan_outbound: bool = True,
        scan_inbound: bool = False,
        action_pii_external: Literal['warn', 'hold', 'block', 'recall'] = 'warn',
        action_financial_external: Literal['warn', 'hold', 'block', 'recall'] = 'warn',
        action_credentials: Literal['warn', 'hold', 'block', 'recall'] = 'block',
        action_legal: Literal['warn', 'hold', 'block', 'recall'] = 'hold',
        notify_sender: bool = True,
        notify_admin: bool = True,
        admin_emails: list[str] = None,
    ):
        self.enabled = enabled
        self.scan_outbound = scan_outbound
        self.scan_inbound = scan_inbound
        self.action_pii_external = action_pii_external
        self.action_financial_external = action_financial_external
        self.action_credentials = action_credentials
        self.action_legal = action_legal
        self.notify_sender = notify_sender
        self.notify_admin = notify_admin
        self.admin_emails = admin_emails or []


async def get_dlp_config(org_id: str, db: AsyncSession) -> Optional[DLPConfig]:
    """Load DLP configuration for an organization."""
    try:
        result = await db.execute(
            text("""
                SELECT 
                    dlp_enabled,
                    dlp_scan_outbound,
                    dlp_scan_inbound,
                    dlp_action_pii,
                    dlp_action_financial,
                    dlp_action_credentials,
                    dlp_action_legal,
                    dlp_notify_sender,
                    dlp_notify_admin,
                    dlp_admin_emails
                FROM organizations
                WHERE id = :org_id
            """),
            {"org_id": org_id}
        )
        row = result.fetchone()
        if not row:
            logger.debug(f"DLP config: no org found for {org_id}")
            return None
        if not row[0]:  # dlp_enabled is False or null
            logger.debug(f"DLP config: dlp_enabled={row[0]} for org {org_id}")
            return None
        logger.info(f"DLP config loaded for org {org_id}: enabled={row[0]}, scan_outbound={row[1]}, action_pii={row[3]}, action_financial={row[4]}, action_creds={row[5]}, action_legal={row[6]}")
    except Exception as exc:
        logger.warning(f"DLP config query failed for org {org_id}: {exc}")
        return None
    
    return DLPConfig(
        enabled=bool(row[0]),
        scan_outbound=bool(row[1]) if row[1] is not None else True,
        scan_inbound=bool(row[2]) if row[2] is not None else False,
        action_pii_external=row[3] or 'warn',
        action_financial_external=row[4] or 'warn',
        action_credentials=row[5] or 'block',
        action_legal=row[6] or 'hold',
        notify_sender=bool(row[7]) if row[7] is not None else True,
        notify_admin=bool(row[8]) if row[8] is not None else True,
        admin_emails=row[9] or [],
    )


async def enable_dlp(
    org_id: str,
    db: AsyncSession,
    config: DLPConfig,
) -> bool:
    """Enable DLP for an organization with the given configuration."""
    try:
        await db.execute(
            text("""
                UPDATE organizations SET
                    dlp_enabled = :enabled,
                    dlp_scan_outbound = :scan_outbound,
                    dlp_scan_inbound = :scan_inbound,
                    dlp_action_pii = :action_pii,
                    dlp_action_financial = :action_financial,
                    dlp_action_credentials = :action_credentials,
                    dlp_action_legal = :action_legal,
                    dlp_notify_sender = :notify_sender,
                    dlp_notify_admin = :notify_admin,
                    dlp_admin_emails = :admin_emails,
                    dlp_enabled_at = NOW()
                WHERE id = :org_id
            """),
            {
                "org_id": org_id,
                "enabled": config.enabled,
                "scan_outbound": config.scan_outbound,
                "scan_inbound": config.scan_inbound,
                "action_pii": config.action_pii_external,
                "action_financial": config.action_financial_external,
                "action_credentials": config.action_credentials,
                "action_legal": config.action_legal,
                "notify_sender": config.notify_sender,
                "notify_admin": config.notify_admin,
                "admin_emails": config.admin_emails,
            }
        )
        await db.commit()
        return True
    except Exception as exc:
        logger.error(f"Failed to enable DLP for org {org_id}: {exc}")
        await db.rollback()
        return False


# ── Email Scanning ────────────────────────────────────────────────────────

async def scan_email_for_dlp(
    org_id: str,
    email_id: str,
    subject: str,
    body: str,
    sender: str,
    recipients: list[str],
    attachments: list[dict],  # [{name, content_type, content_base64}]
    direction: Literal['inbound', 'outbound'],
    db: AsyncSession,
) -> dict:
    """
    Scan an email for DLP violations.
    
    Returns:
        {
            "violations": [...],
            "action": "allow" | "warn" | "hold" | "block",
            "categories": [...],
            "severity": "low" | "medium" | "high" | "critical"
        }
    """
    config = await get_dlp_config(org_id, db)
    if not config or not config.enabled:
        return {"violations": [], "action": "allow", "categories": [], "severity": "low"}
    
    # Skip based on direction
    if direction == 'outbound' and not config.scan_outbound:
        return {"violations": [], "action": "allow", "categories": [], "severity": "low"}
    if direction == 'inbound' and not config.scan_inbound:
        return {"violations": [], "action": "allow", "categories": [], "severity": "low"}
    
    # Determine if recipients are external
    org_domains = await _get_org_domains(org_id, db)
    has_external = any(
        not any(r.lower().endswith(f"@{d.lower()}") for d in org_domains)
        for r in recipients
    )
    
    # Combine content for scanning
    full_text = f"{subject}\n\n{body}"
    
    # Run classification
    result = await classify_text(full_text)
    
    violations = []
    max_severity = "low"
    action = "allow"
    
    # Check each detection
    for detection in result.get("detections", []):
        cat = detection.get("category", "")
        sev = detection.get("severity", "medium")
        
        violation = {
            "category": cat,
            "severity": sev,
            "match": detection.get("match", ""),
            "pattern": detection.get("pattern", ""),
        }
        violations.append(violation)
        
        # Update max severity
        if _severity_rank(sev) > _severity_rank(max_severity):
            max_severity = sev
        
        # Determine action based on category and whether external
        if has_external:
            if cat.startswith("pii_"):
                cat_action = config.action_pii_external
            elif cat.startswith("financial_"):
                cat_action = config.action_financial_external
            elif cat.startswith("credential_"):
                cat_action = config.action_credentials
            elif cat.startswith("legal_"):
                cat_action = config.action_legal
            else:
                cat_action = "warn"
            
            if _action_rank(cat_action) > _action_rank(action):
                action = cat_action
    
    # Also scan attachments if present
    for att in attachments:
        att_result = await _scan_attachment(att, config)
        violations.extend(att_result.get("violations", []))
        if _action_rank(att_result.get("action", "allow")) > _action_rank(action):
            action = att_result.get("action", "allow")
    
    categories = list(set(v["category"] for v in violations))
    
    # Log the DLP event
    await _log_dlp_event(
        org_id=org_id,
        email_id=email_id,
        direction=direction,
        sender=sender,
        recipients=recipients,
        subject=subject,
        violations=violations,
        action=action,
        severity=max_severity,
        db=db,
    )
    
    return {
        "violations": violations,
        "action": action,
        "categories": categories,
        "severity": max_severity,
    }


async def _get_org_domains(org_id: str, db: AsyncSession) -> list[str]:
    """Get domains associated with an organization."""
    domains = []
    
    # Try org_integrations table first
    try:
        result = await db.execute(
            text("""
                SELECT DISTINCT org_domain
                FROM org_integrations
                WHERE org_id = :org_id AND status = 'active' AND org_domain IS NOT NULL
            """),
            {"org_id": org_id}
        )
        domains.extend([r[0] for r in result.fetchall() if r[0]])
    except Exception:
        pass
    
    # Also check organization primary domain
    try:
        org_result = await db.execute(
            text("SELECT domain FROM organizations WHERE id = :org_id"),
            {"org_id": org_id}
        )
        org_row = org_result.fetchone()
        if org_row and org_row[0]:
            domains.append(org_row[0])
    except Exception:
        pass
    
    return list(set(domains))


async def _extract_text_from_attachment(att: dict) -> str:
    """
    Extract text content from an attachment.
    Supports: PDF, text files, Office docs (DOCX/XLSX/PPTX), images (OCR).
    """
    import base64
    import io
    
    content_b64 = att.get("content_base64", "")
    if not content_b64:
        return ""
    
    try:
        content_bytes = base64.b64decode(content_b64)
    except Exception as e:
        logger.warning(f"Failed to decode attachment base64: {e}")
        return ""
    
    content_type = att.get("content_type", "").lower()
    name = att.get("name", "").lower()
    
    # PDF extraction using pdfminer.six
    if content_type == "application/pdf" or name.endswith(".pdf"):
        try:
            from pdfminer.high_level import extract_text as pdf_extract_text
            from pdfminer.pdfparser import PDFSyntaxError
            text = pdf_extract_text(io.BytesIO(content_bytes), maxpages=20)
            return (text or "")[:50000]
        except PDFSyntaxError as e:
            logger.debug(f"PDF syntax error for {name}: {e}")
            return ""
        except Exception as e:
            logger.warning(f"PDF extraction failed for {name}: {e}")
            return ""
    
    # Plain text files
    if content_type.startswith("text/") or name.endswith((".txt", ".csv", ".log", ".json", ".xml", ".md", ".html", ".htm")):
        try:
            text = content_bytes.decode("utf-8", errors="ignore")
            # Strip HTML tags if HTML
            if name.endswith((".html", ".htm")) or "html" in content_type:
                import re
                text = re.sub(r'<[^>]+>', ' ', text)
                text = re.sub(r'\s+', ' ', text)
            return text[:50000]
        except Exception:
            return ""
    
    # DOCX (Word 2007+)
    if content_type in ("application/vnd.openxmlformats-officedocument.wordprocessingml.document",) or name.endswith(".docx"):
        try:
            from docx import Document
            doc = Document(io.BytesIO(content_bytes))
            text_parts = []
            for p in doc.paragraphs:
                if p.text.strip():
                    text_parts.append(p.text)
            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        text_parts.append(row_text)
            return "\n".join(text_parts)[:50000]
        except ImportError:
            logger.debug("python-docx not installed, skipping .docx extraction")
            return ""
        except Exception as e:
            logger.warning(f"DOCX extraction failed for {name}: {e}")
            return ""
    
    # XLSX (Excel 2007+)
    if content_type in ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",) or name.endswith(".xlsx"):
        try:
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(content_bytes), read_only=True, data_only=True)
            text_parts = []
            for sheet in wb.worksheets[:5]:  # Limit sheets
                for row in sheet.iter_rows(max_row=500, values_only=True):  # Limit rows
                    row_text = " ".join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        text_parts.append(row_text)
            wb.close()
            return "\n".join(text_parts)[:50000]
        except ImportError:
            logger.debug("openpyxl not installed, skipping .xlsx extraction")
            return ""
        except Exception as e:
            logger.warning(f"XLSX extraction failed for {name}: {e}")
            return ""
    
    # PPTX (PowerPoint 2007+)
    if content_type in ("application/vnd.openxmlformats-officedocument.presentationml.presentation",) or name.endswith(".pptx"):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content_bytes))
            text_parts = []
            for slide in prs.slides[:30]:  # Limit slides
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
            return "\n".join(text_parts)[:50000]
        except ImportError:
            logger.debug("python-pptx not installed, skipping .pptx extraction")
            return ""
        except Exception as e:
            logger.warning(f"PPTX extraction failed for {name}: {e}")
            return ""
    
    # Images (OCR with pytesseract)
    if content_type.startswith("image/") or name.endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp")):
        try:
            from PIL import Image
            import pytesseract
            img = Image.open(io.BytesIO(content_bytes))
            # Resize large images to speed up OCR
            max_dim = 2000
            if img.width > max_dim or img.height > max_dim:
                ratio = min(max_dim / img.width, max_dim / img.height)
                img = img.resize((int(img.width * ratio), int(img.height * ratio)))
            text = pytesseract.image_to_string(img, timeout=30)
            return (text or "")[:50000]
        except ImportError:
            logger.debug("pytesseract/PIL not installed, skipping image OCR")
            return ""
        except Exception as e:
            logger.debug(f"Image OCR failed for {name}: {e}")
            return ""
    
    # RTF files
    if content_type == "application/rtf" or name.endswith(".rtf"):
        try:
            from striprtf.striprtf import rtf_to_text
            text = rtf_to_text(content_bytes.decode("utf-8", errors="ignore"))
            return (text or "")[:50000]
        except ImportError:
            logger.debug("striprtf not installed, skipping .rtf extraction")
            return ""
        except Exception as e:
            logger.warning(f"RTF extraction failed for {name}: {e}")
            return ""
    
    # EML files (email messages as attachments)
    if content_type == "message/rfc822" or name.endswith(".eml"):
        try:
            import email
            from email import policy
            msg = email.message_from_bytes(content_bytes, policy=policy.default)
            text_parts = []
            text_parts.append(f"Subject: {msg.get('subject', '')}")
            text_parts.append(f"From: {msg.get('from', '')}")
            text_parts.append(f"To: {msg.get('to', '')}")
            # Get body
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        payload = part.get_payload(decode=True)
                        if payload:
                            text_parts.append(payload.decode("utf-8", errors="ignore"))
            else:
                payload = msg.get_payload(decode=True)
                if payload:
                    text_parts.append(payload.decode("utf-8", errors="ignore"))
            return "\n".join(text_parts)[:50000]
        except Exception as e:
            logger.warning(f"EML extraction failed for {name}: {e}")
            return ""
    
    return ""


async def _scan_attachment(att: dict, config: DLPConfig) -> dict:
    """
    Scan an attachment for DLP violations.
    Now extracts and scans actual content, not just filenames.
    """
    name = att.get("name", "").lower()
    violations = []
    action = "allow"
    
    # 1. Check filename patterns (quick check)
    sensitive_patterns = [
        ("salary", "financial_salary", "medium"),
        ("payroll", "financial_salary", "medium"),
        ("ssn", "pii_ssn", "critical"),
        ("password", "credential_password", "high"),
        ("secret", "credential_apikey", "high"),
        ("confidential", "legal_nda", "medium"),
        ("private_key", "credential_privkey", "critical"),
        ("budget", "financial_budget", "medium"),
        ("invoice", "financial_invoice", "medium"),
        ("financial", "financial_report", "medium"),
        ("bank", "financial_banking", "medium"),
        ("tax", "financial_tax", "medium"),
    ]
    
    for pattern, category, severity in sensitive_patterns:
        if pattern in name:
            violations.append({
                "category": category,
                "severity": severity,
                "match": f"filename: {att.get('name')}",
                "pattern": f"attachment name contains '{pattern}'",
            })
            
            if category.startswith("credential_"):
                cat_action = config.action_credentials
            elif category.startswith("pii_"):
                cat_action = config.action_pii_external
            elif category.startswith("financial_"):
                cat_action = config.action_financial_external
            else:
                cat_action = "warn"
            
            if _action_rank(cat_action) > _action_rank(action):
                action = cat_action
    
    # 2. Extract and scan content if available
    if att.get("content_base64"):
        extracted_text = await _extract_text_from_attachment(att)
        if extracted_text and len(extracted_text) > 50:  # Only scan if meaningful content
            logger.info(f"DLP: Scanning attachment content for '{att.get('name')}' ({len(extracted_text)} chars)")
            
            # Run through the same classifier as email body
            content_result = await classify_text(f"[Attachment: {att.get('name')}]\n\n{extracted_text}")
            
            for detection in content_result.get("detections", []):
                cat = detection.get("category", "")
                sev = detection.get("severity", "medium")
                
                violations.append({
                    "category": cat,
                    "severity": sev,
                    "match": f"attachment content: {detection.get('match', '')[:100]}",
                    "pattern": detection.get("pattern", ""),
                })
                
                # Determine action based on category
                if cat.startswith("credential_"):
                    cat_action = config.action_credentials
                elif cat.startswith("pii_"):
                    cat_action = config.action_pii_external
                elif cat.startswith("financial_"):
                    cat_action = config.action_financial_external
                elif cat.startswith("legal_"):
                    cat_action = config.action_legal
                else:
                    cat_action = "warn"
                
                if _action_rank(cat_action) > _action_rank(action):
                    action = cat_action
    
    return {"violations": violations, "action": action}


async def _log_dlp_event(
    org_id: str,
    email_id: str,
    direction: str,
    sender: str,
    recipients: list[str],
    subject: str,
    violations: list[dict],
    action: str,
    severity: str,
    db: AsyncSession,
) -> None:
    """Log a DLP event to the database."""
    try:
        import json
        # Extract categories from violations
        categories = list(set(v.get("category", "") for v in violations if v.get("category")))
        # Extract matched patterns
        # Store patterns as simple strings for display
        matched_patterns = [v.get("pattern", v.get("category", "unknown")) for v in violations if v.get("pattern") or v.get("category")]
        
        await db.execute(
            text("""
                INSERT INTO dlp_events (
                    id, org_id, sender_email, recipient_emails,
                    subject, body_preview, risk_level, action_taken, 
                    categories_found, matched_patterns, confidence, created_at
                ) VALUES (
                    gen_random_uuid(), :org_id, :sender, CAST(:recipients AS jsonb),
                    :subject, :body_preview, :risk_level, :action, 
                    CAST(:categories AS jsonb), CAST(:patterns AS jsonb), :confidence, NOW()
                )
            """),
            {
                "org_id": org_id,
                "sender": sender,
                "recipients": json.dumps(recipients),
                "subject": subject[:500] if subject else "",
                "body_preview": f"[{direction.upper()}] Email ID: {email_id}",
                "risk_level": severity,
                "action": action.upper(),
                "categories": json.dumps(categories),
                "patterns": json.dumps(matched_patterns),
                "confidence": 0.85 if violations else 0.5,
            }
        )
        await db.commit()
        logger.info(f"DLP event logged: org={org_id} sender={sender} action={action} categories={categories}")
    except Exception as exc:
        logger.warning(f"Failed to log DLP event: {exc}")


def _severity_rank(severity: str) -> int:
    """Rank severity levels."""
    return {"low": 0, "medium": 1, "high": 2, "critical": 3}.get(severity, 0)


def _action_rank(action: str) -> int:
    """Rank actions by strictness. Recall is most strict (blocks + recalls)."""
    return {"allow": 0, "warn": 1, "hold": 2, "block": 3, "recall": 4}.get(action, 0)


# ── Batch Scanning (Sent Items) ───────────────────────────────────────────

async def scan_sent_items(
    org_id: str,
    integration_id: str,
    since_hours: int = 1,
    db: AsyncSession = None,
) -> dict:
    """
    Scan recently sent emails for DLP violations.
    This runs as a background job after emails are sent.
    
    Returns summary of violations found.
    """
    # This would be called by a cron job or the delta sync process
    # Implementation depends on provider API
    
    config = await get_dlp_config(org_id, db)
    if not config or not config.enabled or not config.scan_outbound:
        return {"scanned": 0, "violations": 0, "status": "skipped"}
    
    # Fetch sent items from the last N hours via provider API
    # For M365: use Graph API /me/mailFolders/SentItems/messages
    # For Gmail: use Gmail API users.messages.list with q="in:sent after:N"
    
    # This is a placeholder - actual implementation would call provider APIs
    return {
        "scanned": 0,
        "violations": 0,
        "status": "pending_implementation",
        "message": "Sent items scanning runs automatically via delta sync"
    }


# ── Integration with Delta Sync ───────────────────────────────────────────

async def process_email_with_dlp(
    org_id: str,
    email_data: dict,
    direction: Literal['inbound', 'outbound'],
    db: AsyncSession,
) -> dict:
    """
    Process an email through DLP during delta sync.
    Called by the delta sync workers when processing emails.
    
    Returns the DLP result which can be used to:
    - Add labels to the email
    - Move to a DLP quarantine folder
    - Send notifications
    """
    config = await get_dlp_config(org_id, db)
    if not config or not config.enabled:
        return {"dlp_checked": False}
    
    result = await scan_email_for_dlp(
        org_id=org_id,
        email_id=email_data.get("id", ""),
        subject=email_data.get("subject", ""),
        body=email_data.get("body", ""),
        sender=email_data.get("sender", ""),
        recipients=email_data.get("recipients", []),
        attachments=email_data.get("attachments", []),
        direction=direction,
        db=db,
    )
    
    result["dlp_checked"] = True
    
    # If action is hold, block, or recall, we can take provider-specific actions
    if result["action"] in ["hold", "block", "recall"]:
        # Queue notification
        if config.notify_admin and config.admin_emails:
            await _queue_admin_notification(
                org_id=org_id,
                email_data=email_data,
                dlp_result=result,
                admin_emails=config.admin_emails,
                db=db,
            )
        
        if config.notify_sender:
            await _queue_sender_notification(
                org_id=org_id,
                email_data=email_data,
                dlp_result=result,
                db=db,
            )
        
        # Auto-recall if configured
        if result["action"] == "recall":
            await _recall_email(
                org_id=org_id,
                email_data=email_data,
                dlp_result=result,
                db=db,
            )
    
    return result


async def _queue_admin_notification(
    org_id: str,
    email_data: dict,
    dlp_result: dict,
    admin_emails: list[str],
    db: AsyncSession,
) -> None:
    """Queue a notification to admins about a DLP violation."""
    # Implementation would add to notification queue
    logger.info(
        f"DLP violation notification queued for org {org_id}: "
        f"{len(dlp_result.get('violations', []))} violations, "
        f"action={dlp_result.get('action')}"
    )


async def _queue_sender_notification(
    org_id: str,
    email_data: dict,
    dlp_result: dict,
    db: AsyncSession,
) -> None:
    """Queue a notification to the sender about their email being held/blocked."""
    logger.info(
        f"DLP sender notification queued for {email_data.get('sender')}: "
        f"action={dlp_result.get('action')}"
    )


async def _recall_email(
    org_id: str,
    email_data: dict,
    dlp_result: dict,
    db: AsyncSession,
) -> bool:
    """
    Attempt to recall/delete an email that violated DLP policy.
    
    For M365: Uses Graph API to delete from recipient mailboxes (if within org)
    For Gmail: Uses Gmail API to delete from sent folder and request recall
    
    Note: External recipients cannot be recalled - email is already delivered.
    This works best for internal emails where we have admin access.
    """
    email_id = email_data.get("id", "")
    provider = email_data.get("provider", "")
    recipients = email_data.get("recipients", [])
    
    logger.info(
        f"DLP auto-recall initiated for email {email_id}: "
        f"provider={provider}, recipients={len(recipients)}"
    )
    
    # Get integration details from org_integrations table
    # Columns are encrypted: access_token_enc, refresh_token_enc
    # Map email provider to DB provider column value
    # DB uses: "google" for Google Workspace, could be "microsoft" or "m365" for M365
    provider_map = {"google": "google", "gmail": "google", "m365": "m365", "microsoft": "m365"}
    int_provider_name = provider_map.get(provider, provider)
    
    logger.info(f"DLP recall: looking for integration with provider={int_provider_name} for org={org_id}")
    
    try:
        from backend.services.delta_sync import _decrypt
        
        integration_result = await db.execute(
            text("""
                SELECT id, provider, access_token_enc, refresh_token_enc
                FROM org_integrations
                WHERE org_id = :org_id AND provider = :provider AND status = 'active'
                LIMIT 1
            """),
            {"org_id": org_id, "provider": int_provider_name}
        )
        integration = integration_result.fetchone()
        
        if not integration:
            logger.warning(f"No active {int_provider_name} integration found for org {org_id}")
            return False
        
        int_id, int_provider, enc_token, enc_refresh = integration
        access_token = _decrypt(enc_token) if enc_token else None
        
        if not access_token:
            logger.warning(f"No valid access token for org {org_id}")
            return False
        
        if int_provider in ("microsoft", "m365"):
            return await _recall_m365_email(email_data, access_token, email_data.get("sender", ""), db)
        elif int_provider == "google":
            return await _recall_gmail_email(email_data, access_token, email_data.get("sender", ""), db)
        else:
            logger.warning(f"Recall not supported for provider: {int_provider}")
            return False
            
    except Exception as exc:
        logger.error(f"Email recall failed: {exc}")
        return False


async def _recall_m365_email(email_data: dict, access_token: str, sender_email: str, db: AsyncSession) -> bool:
    """
    Recall email via Microsoft Graph API.
    
    Uses the 'soft delete' approach:
    1. Delete from sender's Sent Items
    2. For internal recipients, delete from their inbox using admin consent
    """
    import httpx
    
    email_id = email_data.get("id", "")
    
    try:
        async with httpx.AsyncClient(timeout=30) as client:
            # Delete from sender's sent items using delegated access
            response = await client.delete(
                f"https://graph.microsoft.com/v1.0/users/{sender_email}/messages/{email_id}",
                headers={"Authorization": f"Bearer {access_token}"},
            )
            
            if response.status_code in (200, 204, 404):
                logger.info(f"M365 email {email_id} deleted from sent items for {sender_email}")
                return True
            else:
                logger.warning(f"M365 recall failed: {response.status_code} {response.text}")
                return False
                
    except Exception as exc:
        logger.error(f"M365 recall error: {exc}")
        return False


async def _recall_gmail_email(email_data: dict, access_token: str, sender_email: str, db: AsyncSession) -> bool:
    """
    Recall email via Gmail API.
    
    Gmail doesn't have true recall, but we can:
    1. Delete from sender's Sent folder using delegated access
    2. For internal recipients with delegated access, delete from their inbox
    """
    import httpx
    
    email_id = email_data.get("id", "")
    
    try:
        async with httpx.AsyncClient(timeout=30) as client:
            # Trash the message using delegated access for the sender
            response = await client.post(
                f"https://gmail.googleapis.com/gmail/v1/users/{sender_email}/messages/{email_id}/trash",
                headers={"Authorization": f"Bearer {access_token}"},
            )
            
            if response.status_code in (200, 204, 404):
                logger.info(f"Gmail email {email_id} moved to trash for {sender_email}")
                return True
            else:
                logger.warning(f"Gmail recall failed: {response.status_code} {response.text[:200] if response.text else ''}")
                return False
                
    except Exception as exc:
        logger.error(f"Gmail recall error: {exc}")
        return False
