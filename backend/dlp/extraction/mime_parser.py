"""Bounded MIME parsing with isolated attachment extraction."""

from __future__ import annotations

import asyncio
import hashlib
import io
import multiprocessing
import queue
import zipfile
from dataclasses import dataclass, field
from email import policy
from email.message import Message
from email.parser import BytesParser
from pathlib import Path
from typing import Any

from bs4 import BeautifulSoup


@dataclass(frozen=True)
class MimeExtractionLimits:
    max_mime_bytes: int = 25 * 1024 * 1024
    max_parts: int = 100
    max_part_bytes: int = 10 * 1024 * 1024
    max_text_bytes: int = 2 * 1024 * 1024
    max_archive_uncompressed_bytes: int = 50 * 1024 * 1024
    max_archive_ratio: int = 100
    attachment_timeout_seconds: float = 10.0
    attachment_memory_bytes: int = 256 * 1024 * 1024


@dataclass(frozen=True)
class ExtractionLimitation:
    code: str
    detail: str
    part_index: int | None = None
    filename: str | None = None
    fatal: bool = False


@dataclass(frozen=True)
class ExtractedPart:
    part_index: int
    content_type: str
    disposition: str | None
    filename: str | None
    size_bytes: int
    sha256: str
    text: str = ""
    limitation: ExtractionLimitation | None = None


@dataclass(frozen=True)
class MimeExtractionResult:
    subject: str
    sender: str
    recipients: tuple[str, ...]
    text: str
    parts: tuple[ExtractedPart, ...]
    limitations: tuple[ExtractionLimitation, ...]

    @property
    def is_complete(self) -> bool:
        return not self.limitations


class MimeExtractionError(ValueError):
    pass


class SafeMimeExtractor:
    def __init__(
        self, limits: MimeExtractionLimits | None = None
    ) -> None:
        self.limits = limits or MimeExtractionLimits()

    async def extract(self, raw_mime: bytes) -> MimeExtractionResult:
        if len(raw_mime) > self.limits.max_mime_bytes:
            raise MimeExtractionError(
                f"MIME exceeds {self.limits.max_mime_bytes} bytes"
            )
        try:
            message = BytesParser(policy=policy.default).parsebytes(raw_mime)
        except Exception as exc:
            raise MimeExtractionError(f"MIME parsing failed: {exc}") from exc

        limitations: list[ExtractionLimitation] = []
        if message.defects:
            limitations.append(
                ExtractionLimitation(
                    code="malformed_mime",
                    detail="; ".join(
                        type(defect).__name__ for defect in message.defects
                    ),
                )
            )

        parts: list[ExtractedPart] = []
        text_sections: list[str] = []
        leaf_parts = list(_leaf_parts(message))
        if len(leaf_parts) > self.limits.max_parts:
            limitations.append(
                ExtractionLimitation(
                    code="part_limit_exceeded",
                    detail=(
                        f"Only the first {self.limits.max_parts} MIME "
                        "parts were inspected"
                    ),
                    fatal=True,
                )
            )
            leaf_parts = leaf_parts[: self.limits.max_parts]

        for part_index, part in enumerate(leaf_parts):
            extracted = await self._extract_part(part, part_index)
            parts.append(extracted)
            if extracted.limitation is not None:
                limitations.append(extracted.limitation)
            if extracted.text:
                text_sections.append(extracted.text)

        header_text = "\n".join(
            (
                f"Subject: {_bounded_header(message.get('subject'))}",
                f"From: {_bounded_header(message.get('from'))}",
                f"To: {_bounded_header(message.get('to'))}",
            )
        )
        combined, was_truncated = _truncate_utf8(
            "\n\n".join([header_text, *text_sections]),
            self.limits.max_text_bytes,
        )
        if was_truncated:
            limitations.append(
                ExtractionLimitation(
                    code="text_limit_exceeded",
                    detail=(
                        "Classifier text was truncated at "
                        f"{self.limits.max_text_bytes} UTF-8 bytes"
                    ),
                    fatal=True,
                )
            )

        return MimeExtractionResult(
            subject=_bounded_header(message.get("subject")),
            sender=_bounded_header(message.get("from")),
            recipients=tuple(
                _bounded_header(value)
                for value in message.get_all("to", [])
            ),
            text=combined,
            parts=tuple(parts),
            limitations=tuple(limitations),
        )

    async def _extract_part(
        self, part: Message, part_index: int
    ) -> ExtractedPart:
        content_type = part.get_content_type().lower()
        disposition = part.get_content_disposition()
        filename = _bounded_header(part.get_filename()) or None
        payload = part.get_payload(decode=True) or b""
        digest = hashlib.sha256(payload).hexdigest()

        if len(payload) > self.limits.max_part_bytes:
            limitation = ExtractionLimitation(
                code="part_size_limit_exceeded",
                detail=(
                    f"Part exceeds {self.limits.max_part_bytes} bytes"
                ),
                part_index=part_index,
                filename=filename,
                fatal=True,
            )
            return ExtractedPart(
                part_index,
                content_type,
                disposition,
                filename,
                len(payload),
                digest,
                limitation=limitation,
            )

        if _is_encrypted_part(content_type, payload):
            limitation = ExtractionLimitation(
                code="encrypted_content",
                detail="Encrypted content cannot be inspected",
                part_index=part_index,
                filename=filename,
                fatal=True,
            )
            return ExtractedPart(
                part_index,
                content_type,
                disposition,
                filename,
                len(payload),
                digest,
                limitation=limitation,
            )

        if content_type == "text/plain":
            text = _decode_text(payload, part.get_content_charset())
            return ExtractedPart(
                part_index,
                content_type,
                disposition,
                filename,
                len(payload),
                digest,
                text=text,
            )
        if content_type == "text/html":
            html = _decode_text(payload, part.get_content_charset())
            text = BeautifulSoup(html, "html.parser").get_text(
                separator=" ", strip=True
            )
            return ExtractedPart(
                part_index,
                content_type,
                disposition,
                filename,
                len(payload),
                digest,
                text=text,
            )
        if content_type in {"text/csv", "application/csv"}:
            text = _decode_text(payload, part.get_content_charset())
            return ExtractedPart(
                part_index,
                content_type,
                disposition,
                filename,
                len(payload),
                digest,
                text=text,
            )

        extractor_type = _attachment_extractor_type(
            content_type, filename
        )
        if extractor_type is None:
            limitation = ExtractionLimitation(
                code="unsupported_content_type",
                detail=f"No extractor for {content_type}",
                part_index=part_index,
                filename=filename,
            )
            return ExtractedPart(
                part_index,
                content_type,
                disposition,
                filename,
                len(payload),
                digest,
                limitation=limitation,
            )

        archive_limitation = _validate_archive(
            payload, self.limits, part_index, filename
        )
        if archive_limitation is not None:
            return ExtractedPart(
                part_index,
                content_type,
                disposition,
                filename,
                len(payload),
                digest,
                limitation=archive_limitation,
            )

        text, limitation = await asyncio.to_thread(
            _run_attachment_process,
            extractor_type,
            payload,
            self.limits,
            part_index,
            filename,
        )
        return ExtractedPart(
            part_index,
            content_type,
            disposition,
            filename,
            len(payload),
            digest,
            text=text,
            limitation=limitation,
        )


def _leaf_parts(message: Message):
    if message.is_multipart():
        for part in message.iter_parts():
            yield from _leaf_parts(part)
    else:
        yield message


def _bounded_header(value: Any, limit: int = 4096) -> str:
    return str(value or "").replace("\x00", "")[:limit]


def _decode_text(payload: bytes, charset: str | None) -> str:
    try:
        return payload.decode(charset or "utf-8", errors="replace")
    except LookupError:
        return payload.decode("utf-8", errors="replace")


def _truncate_utf8(value: str, max_bytes: int) -> tuple[str, bool]:
    encoded = value.encode("utf-8")
    if len(encoded) <= max_bytes:
        return value, False
    return (
        encoded[:max_bytes].decode("utf-8", errors="ignore"),
        True,
    )


def _is_encrypted_part(content_type: str, payload: bytes) -> bool:
    return (
        content_type
        in {
            "application/pkcs7-mime",
            "application/x-pkcs7-mime",
            "application/pgp-encrypted",
        }
        or (
            content_type == "application/pdf"
            and b"/Encrypt" in payload[:1_000_000]
        )
    )


def _attachment_extractor_type(
    content_type: str, filename: str | None
) -> str | None:
    by_content_type = {
        "application/pdf": "pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
        "application/rtf": "rtf",
        "text/rtf": "rtf",
    }
    if content_type in by_content_type:
        return by_content_type[content_type]
    extension = Path(filename or "").suffix.lower().lstrip(".")
    if content_type == "application/octet-stream" and extension in {
        "pdf",
        "docx",
        "xlsx",
        "pptx",
        "rtf",
    }:
        return extension
    return None


def _validate_archive(
    payload: bytes,
    limits: MimeExtractionLimits,
    part_index: int,
    filename: str | None,
) -> ExtractionLimitation | None:
    if not zipfile.is_zipfile(io.BytesIO(payload)):
        return None
    try:
        with zipfile.ZipFile(io.BytesIO(payload)) as archive:
            total_uncompressed = sum(
                member.file_size for member in archive.infolist()
            )
            total_compressed = max(
                sum(member.compress_size for member in archive.infolist()),
                1,
            )
    except zipfile.BadZipFile:
        return ExtractionLimitation(
            code="malformed_archive",
            detail="Attachment archive is malformed",
            part_index=part_index,
            filename=filename,
            fatal=True,
        )
    if (
        total_uncompressed > limits.max_archive_uncompressed_bytes
        or total_uncompressed / total_compressed
        > limits.max_archive_ratio
    ):
        return ExtractionLimitation(
            code="archive_expansion_limit",
            detail="Attachment archive exceeds safe expansion limits",
            part_index=part_index,
            filename=filename,
            fatal=True,
        )
    return None


def _run_attachment_process(
    extractor_type: str,
    payload: bytes,
    limits: MimeExtractionLimits,
    part_index: int,
    filename: str | None,
) -> tuple[str, ExtractionLimitation | None]:
    context = multiprocessing.get_context("spawn")
    result_queue = context.Queue(maxsize=1)
    process = context.Process(
        target=_attachment_worker,
        args=(
            extractor_type,
            payload,
            limits.max_text_bytes,
            limits.attachment_memory_bytes,
            result_queue,
        ),
        daemon=True,
    )
    process.start()
    process.join(limits.attachment_timeout_seconds)
    if process.is_alive():
        process.terminate()
        process.join(2)
        return "", ExtractionLimitation(
            code="extractor_timeout",
            detail="Attachment extraction timed out",
            part_index=part_index,
            filename=filename,
            fatal=True,
        )
    try:
        status, value = result_queue.get(timeout=1)
    except queue.Empty:
        status, value = (
            "error",
            f"Extractor exited with code {process.exitcode}",
        )
    finally:
        result_queue.close()
    if status == "ok":
        text, truncated = _truncate_utf8(
            str(value), limits.max_text_bytes
        )
        limitation = (
            ExtractionLimitation(
                code="attachment_text_limit",
                detail="Attachment text was truncated",
                part_index=part_index,
                filename=filename,
                fatal=True,
            )
            if truncated
            else None
        )
        return text, limitation
    return "", ExtractionLimitation(
        code="extractor_error",
        detail=str(value)[:1000],
        part_index=part_index,
        filename=filename,
        fatal=True,
    )


def _attachment_worker(
    extractor_type: str,
    payload: bytes,
    max_text_bytes: int,
    memory_bytes: int,
    result_queue,
) -> None:
    _apply_resource_limits(memory_bytes)
    try:
        text = _extract_attachment_text(
            extractor_type, payload, max_text_bytes
        )
        result_queue.put(("ok", text))
    except Exception as exc:
        result_queue.put(
            ("error", f"{type(exc).__name__}: {exc}")
        )


def _apply_resource_limits(memory_bytes: int) -> None:
    try:
        import resource

        resource.setrlimit(
            resource.RLIMIT_AS, (memory_bytes, memory_bytes)
        )
        resource.setrlimit(resource.RLIMIT_CPU, (8, 8))
        resource.setrlimit(resource.RLIMIT_NOFILE, (32, 32))
    except (ImportError, OSError, ValueError):
        # Windows lacks resource limits; the timeout still bounds execution.
        return


def _extract_attachment_text(
    extractor_type: str, payload: bytes, max_text_bytes: int
) -> str:
    source = io.BytesIO(payload)
    if extractor_type == "pdf":
        from pdfminer.high_level import extract_text

        return extract_text(source)
    if extractor_type == "docx":
        from docx import Document

        document = Document(source)
        return "\n".join(
            paragraph.text for paragraph in document.paragraphs
        )
    if extractor_type == "xlsx":
        from openpyxl import load_workbook

        workbook = load_workbook(
            source, read_only=True, data_only=True
        )
        values: list[str] = []
        current_bytes = 0
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                line = "\t".join(
                    "" if value is None else str(value) for value in row
                )
                current_bytes += len(line.encode("utf-8"))
                values.append(line)
                if current_bytes >= max_text_bytes:
                    return "\n".join(values)
        return "\n".join(values)
    if extractor_type == "pptx":
        from pptx import Presentation

        presentation = Presentation(source)
        return "\n".join(
            shape.text
            for slide in presentation.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )
    if extractor_type == "rtf":
        from striprtf.striprtf import rtf_to_text

        return rtf_to_text(payload.decode("utf-8", errors="replace"))
    raise ValueError(f"Unsupported extractor type: {extractor_type}")
